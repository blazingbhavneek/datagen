import asyncio
import json
import logging
import os
import pickle
import re
import subprocess
import tempfile
import time
import uuid
from dataclasses import dataclass
from parser.configs import DocParserConfig
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import chromadb
import docx
import html2text
import numpy as np
import pandas as pd
from bs4 import BeautifulSoup
from chromadb.config import Settings
from openai import OpenAI
from pptx import Presentation
from tqdm.asyncio import tqdm


class DocumentConverter:
    """Converts various document formats to markdown."""

    def __init__(self, config: DocParserConfig):
        self.config = config

    def convert_to_markdown(self, file_path: str) -> str:
        """
        Convert single file to markdown

        Returns:
            - Markdown string with metadata header
        """
        file_ext = Path(file_path).suffix.lower()

        if file_ext == ".pdf":
            return self._convert_pdf(file_path)
        elif file_ext == ".docx":
            return self._convert_docx(file_path)
        elif file_ext in [".html", ".htm"]:
            return self._convert_html(file_path)
        elif file_ext == ".md":
            return self._convert_md(file_path)
        elif file_ext == ".txt":
            return self._convert_txt(file_path)
        elif file_ext == ".xlsx":
            return self._convert_xlsx(file_path)
        elif file_ext == ".pptx":
            return self._convert_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

    def _convert_pdf(self, path: str) -> str:
        """Use MinerU for PDF extraction."""
        # Create a temporary directory for MinerU output
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Run MinerU CLI command
                cmd = [
                    "mineru",
                    "-p",
                    path,
                    "-o",
                    temp_dir,
                ]
                subprocess.run(cmd, check=True)

                # MinerU typically outputs to a subdirectory named after the file
                # Look for markdown files in the output directory
                output_files = list(Path(temp_dir).rglob("*.md"))

                if not output_files:
                    raise RuntimeError(f"No markdown output found for {path}")

                # Read the markdown content
                md_content = output_files[0].read_text(encoding="utf-8")
                return md_content

            except subprocess.CalledProcessError as e:
                raise RuntimeError(f"MinerU failed to process {path}: {e.stderr}")
            except Exception as e:
                raise RuntimeError(f"Failed to convert PDF {path}: {e}")

    def _convert_docx(self, path: str) -> str:
        """Use python-docx or mammoth."""
        if docx is None:
            raise ImportError("python-docx library is required for DOCX conversion")

        try:
            doc = docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise RuntimeError(f"Failed to convert DOCX {path}: {e}")

    def _convert_html(self, path: str) -> str:
        """Use BeautifulSoup + html2text."""
        if BeautifulSoup is None or html2text is None:
            raise ImportError(
                "beautifulsoup4 and html2text libraries are required for HTML conversion"
            )

        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()

            soup = BeautifulSoup(content, "html.parser")
            h = html2text.HTML2Text()
            h.ignore_links = True
            h.body_width = 0  # Don't wrap lines
            return h.handle(str(soup))
        except Exception as e:
            raise RuntimeError(f"Failed to convert HTML {path}: {e}")

    def _convert_md(self, path: str) -> str:
        """Read markdown file directly."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Failed to read MD file {path}: {e}")

    def _convert_txt(self, path: str) -> str:
        """Read text file directly."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Failed to read TXT file {path}: {e}")

    def _convert_xlsx(self, path: str) -> str:
        """Use openpyxl/pandas to convert tables."""
        if pd is None:
            raise ImportError("pandas library is required for XLSX conversion")

        try:
            # Read all sheets and convert to markdown tables
            xl_file = pd.ExcelFile(path)
            sheets = []
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                # Convert DataFrame to markdown table
                table_md = df.to_markdown(index=False)
                sheets.append(f"## Sheet: {sheet_name}\n\n{table_md}\n")

            return "\n".join(sheets)
        except Exception as e:
            raise RuntimeError(f"Failed to convert XLSX {path}: {e}")

    def _convert_pptx(self, path: str) -> str:
        """Use python-pptx to extract text/notes."""
        if Presentation is None:
            raise ImportError("python-pptx library is required for PPTX conversion")

        try:
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_content = [f"### Slide {i+1}"]

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                # Extract notes
                if slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        slide_content.append(f"\nNotes:\n{notes}")

                slides.append("\n".join(slide_content))

            return "\n\n".join(slides)
        except Exception as e:
            raise RuntimeError(f"Failed to convert PPTX {path}: {e}")
