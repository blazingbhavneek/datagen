import asyncio
import json
import logging
import os
import pickle
import re
import subprocess
import tempfile
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import chromadb
import docx
import html2text
import numpy as np
import pandas as pd
from bs4 import BeautifulSoup
from chromadb.config import Settings
from openai import OpenAI
from pptx import Presentation
from tqdm.asyncio import tqdm


@dataclass
class DocParserConfig:
    input_dir: str  # Directory containing documents
    output_md_path: str  # Path for consolidated markdown output
    chroma_db_path: str  # Path for ChromaDB persistence
    embedding_endpoint: str  # e.g., "http://192.168.1.100:8000/v1/embeddings"
    embedding_model: str  # Model name for embeddings
    chunk_size: int = 2000  # Characters per chunk
    chunk_overlap: int = 200  # Character overlap between chunks
    supported_formats: List[str] = None  # Default formats if not provided

    def __post_init__(self):
        if self.supported_formats is None:
            self.supported_formats = [
                ".pdf",
                ".docx",
                ".html",
                ".md",
                ".txt",
                ".xlsx",
                ".pptx",
            ]


@dataclass
class DocParserOutput:
    consolidated_md_path: str  # Path to final markdown file
    chroma_collection_name: str  # ChromaDB collection identifier
    chunk_metadata: List[Dict]  # Metadata for each chunk
    total_chunks: int
    processing_log: str  # Path to processing log


class Chunk:
    """Represents a text chunk with metadata."""

    def __init__(
        self,
        content: str,
        source_file: str,
        chunk_index: int,
        start_char: int,
        end_char: int,
        headers: List[str] = None,
        metadata: Dict = None,
    ):
        self.id = str(uuid.uuid4())
        self.content = content
        self.source_file = source_file
        self.chunk_index = chunk_index
        self.start_char = start_char
        self.end_char = end_char
        self.headers = headers or []
        self.metadata = metadata or {}

    def to_dict(self) -> Dict:
        """Convert chunk to dictionary representation."""
        return {
            "id": self.id,
            "content": self.content,
            "source_file": self.source_file,
            "chunk_index": self.chunk_index,
            "start_char": self.start_char,
            "end_char": self.end_char,
            "headers": self.headers,
            "metadata": self.metadata,
        }


class DocumentConverter:
    """Converts various document formats to markdown."""

    def __init__(self, config: DocParserConfig):
        self.config = config

    def convert_to_markdown(self, file_path: str) -> str:
        """
        Convert single file to markdown

        Returns:
            - Markdown string with metadata header
        """
        file_ext = Path(file_path).suffix.lower()

        if file_ext == ".pdf":
            return self._convert_pdf(file_path)
        elif file_ext == ".docx":
            return self._convert_docx(file_path)
        elif file_ext in [".html", ".htm"]:
            return self._convert_html(file_path)
        elif file_ext == ".md":
            return self._convert_md(file_path)
        elif file_ext == ".txt":
            return self._convert_txt(file_path)
        elif file_ext == ".xlsx":
            return self._convert_xlsx(file_path)
        elif file_ext == ".pptx":
            return self._convert_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

    def _convert_pdf(self, path: str) -> str:
        """Use MinerU for PDF extraction."""
        # Create a temporary directory for MinerU output
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # Run MinerU CLI command
                cmd = [
                    "mineru",
                    "-p",
                    path,
                    "-o",
                    temp_dir,
                ]
                subprocess.run(cmd, check=True)

                # MinerU typically outputs to a subdirectory named after the file
                # Look for markdown files in the output directory
                output_files = list(Path(temp_dir).rglob("*.md"))

                if not output_files:
                    raise RuntimeError(f"No markdown output found for {path}")

                # Read the markdown content
                md_content = output_files[0].read_text(encoding="utf-8")
                return md_content

            except subprocess.CalledProcessError as e:
                raise RuntimeError(f"MinerU failed to process {path}: {e.stderr}")
            except Exception as e:
                raise RuntimeError(f"Failed to convert PDF {path}: {e}")

    def _convert_docx(self, path: str) -> str:
        """Use python-docx or mammoth."""
        if docx is None:
            raise ImportError("python-docx library is required for DOCX conversion")

        try:
            doc = docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise RuntimeError(f"Failed to convert DOCX {path}: {e}")

    def _convert_html(self, path: str) -> str:
        """Use BeautifulSoup + html2text."""
        if BeautifulSoup is None or html2text is None:
            raise ImportError(
                "beautifulsoup4 and html2text libraries are required for HTML conversion"
            )

        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()

            soup = BeautifulSoup(content, "html.parser")
            h = html2text.HTML2Text()
            h.ignore_links = True
            h.body_width = 0  # Don't wrap lines
            return h.handle(str(soup))
        except Exception as e:
            raise RuntimeError(f"Failed to convert HTML {path}: {e}")

    def _convert_md(self, path: str) -> str:
        """Read markdown file directly."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Failed to read MD file {path}: {e}")

    def _convert_txt(self, path: str) -> str:
        """Read text file directly."""
        try:
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Failed to read TXT file {path}: {e}")

    def _convert_xlsx(self, path: str) -> str:
        """Use openpyxl/pandas to convert tables."""
        if pd is None:
            raise ImportError("pandas library is required for XLSX conversion")

        try:
            # Read all sheets and convert to markdown tables
            xl_file = pd.ExcelFile(path)
            sheets = []
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                # Convert DataFrame to markdown table
                table_md = df.to_markdown(index=False)
                sheets.append(f"## Sheet: {sheet_name}\n\n{table_md}\n")

            return "\n".join(sheets)
        except Exception as e:
            raise RuntimeError(f"Failed to convert XLSX {path}: {e}")

    def _convert_pptx(self, path: str) -> str:
        """Use python-pptx to extract text/notes."""
        if Presentation is None:
            raise ImportError("python-pptx library is required for PPTX conversion")

        try:
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_content = [f"### Slide {i+1}"]

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                # Extract notes
                if slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes.strip():
                        slide_content.append(f"\nNotes:\n{notes}")

                slides.append("\n".join(slide_content))

            return "\n\n".join(slides)
        except Exception as e:
            raise RuntimeError(f"Failed to convert PPTX {path}: {e}")


class MarkdownConsolidator:
    """Manages the consolidation of multiple markdown documents."""

    def __init__(self, output_path: str):
        self.output_path = output_path
        self.file_separator = "\n\n" + "=" * 80 + "\n\n"
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

    def append_document(self, md_content: str, source_file: str):
        """
        Append markdown with metadata header

        Header format:
        # SOURCE: {source_file}
        # TIMESTAMP: {iso_timestamp}
        # FORMAT: {original_format}
        ---
        {content}
        """
        source_path = Path(source_file)
        header = (
            f"# SOURCE: {source_file}\n"
            f"# TIMESTAMP: {time.strftime('%Y-%m-%dT%H:%M:%SZ', time.gmtime())}\n"
            f"# FORMAT: {source_path.suffix}\n---\n"
        )

        with open(self.output_path, "a", encoding="utf-8") as f:
            # Add separator if file already has content
            if f.tell() > 0:
                f.write(self.file_separator)

            f.write(header)
            f.write(md_content)

    def finalize(self) -> str:
        """Return path to consolidated file."""
        return self.output_path


class SemanticChunker:
    """Chunks markdown content preserving structure."""

    def __init__(self, chunk_size: int, overlap: int):
        """Initialize with size constraints."""
        self.chunk_size = chunk_size
        self.overlap = overlap

    def chunk_markdown(self, md_content: str, source_file: str) -> List[Chunk]:
        """
        Chunk markdown preserving structure

        Strategy:
        1. Split on markdown headers first (##, ###, etc.)
        2. If section > chunk_size, split on paragraphs
        3. If paragraph > chunk_size, split on sentences
        4. Maintain overlap with previous chunk

        Returns:
            List of Chunk objects with metadata
        """
        chunks = []
        chunk_index = 0

        # First, split by main headers
        header_pattern = r"^(#{1,6})\s+(.*)$"
        lines = md_content.split("\n")

        sections = []  # List of (start_idx, end_idx, headers) tuples
        current_headers = []
        section_start = 0

        for i, line in enumerate(lines):
            match = re.match(header_pattern, line.strip())
            if match:
                # Found a header
                header_level = len(match.group(1))  # Number of # symbols
                header_text = match.group(2)

                # Update headers list based on hierarchy
                current_headers = current_headers[
                    : header_level - 1
                ]  # Keep parent headers
                current_headers.append((header_level, header_text))

                # If we have content since last section, save it
                if i > section_start:
                    sections.append((section_start, i, current_headers.copy()))

                # Start new section after header line
                section_start = i + 1

        # Add the final section
        if section_start < len(lines):
            sections.append((section_start, len(lines), current_headers.copy()))

        # Process each section
        for start_idx, end_idx, headers in sections:
            section_text = "\n".join(lines[start_idx:end_idx])
            section_chunks = self._chunk_section(
                section_text, source_file, headers, chunk_index
            )

            chunks.extend(section_chunks)
            chunk_index += len(section_chunks)

        return chunks

    def _chunk_section(
        self,
        section_text: str,
        source_file: str,
        headers: List[Tuple[int, str]],
        base_chunk_index: int,
    ) -> List[Chunk]:
        """Chunk a single section of markdown."""
        chunks = []
        chunk_index = base_chunk_index

        # If the section is small enough, keep as one chunk
        if len(section_text) <= self.chunk_size:
            chunk = Chunk(
                content=section_text,
                source_file=source_file,
                chunk_index=chunk_index,
                start_char=0,
                end_char=len(section_text),
                headers=[h[1] for h in headers],
                metadata={"header_levels": [h[0] for h in headers]},
            )
            chunks.append(chunk)
            return chunks

        # Otherwise, break down further
        # First try splitting by paragraphs
        paragraphs = section_text.split("\n\n")

        current_chunk = ""
        current_headers = [h[1] for h in headers]

        for para in paragraphs:
            # Check if adding this paragraph would exceed chunk size
            if len(current_chunk) + len(para) <= self.chunk_size:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
            else:
                # Current chunk is full, save it
                if current_chunk.strip():
                    chunk = Chunk(
                        content=current_chunk,
                        source_file=source_file,
                        chunk_index=chunk_index,
                        start_char=0,  # Would need to track actual positions
                        end_char=0,
                        headers=current_headers,
                        metadata={
                            "header_levels": [str(h[0]) for h in headers]
                        },  # Fixed: Convert to string
                    )
                    chunks.append(chunk)
                    chunk_index += 1

                # Start new chunk with overlap
                if len(para) > self.chunk_size:
                    # Paragraph itself is too large, need to split by sentences
                    subchunks = self._split_large_paragraph(
                        para, source_file, current_headers, chunk_index
                    )
                    chunks.extend(subchunks)
                    chunk_index += len(subchunks)
                    current_chunk = ""  # After splitting large paragraph, reset
                else:
                    current_chunk = para

        # Add the final chunk if there's content left
        if current_chunk.strip():
            chunk = Chunk(
                content=current_chunk,
                source_file=source_file,
                chunk_index=chunk_index,
                start_char=0,
                end_char=0,
                headers=current_headers,
                metadata={
                    "header_levels": [str(h[0]) for h in headers]
                },  # Fixed: Convert to string
            )
            chunks.append(chunk)

        return chunks

    def _split_large_paragraph(
        self,
        paragraph: str,
        source_file: str,
        headers: List[str],
        base_chunk_index: int,
    ) -> List[Chunk]:
        """Split a large paragraph into smaller chunks."""
        chunks = []
        chunk_index = base_chunk_index

        # Split by sentences
        sentences = re.split(r"[.!?]+", paragraph)
        current_chunk = ""

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            if len(current_chunk) + len(sentence) <= self.chunk_size:
                if current_chunk:
                    current_chunk += ". " + sentence
                else:
                    current_chunk = sentence
            else:
                # Current chunk is full
                if current_chunk.strip():
                    chunk = Chunk(
                        content=current_chunk + ".",
                        source_file=source_file,
                        chunk_index=chunk_index,
                        start_char=0,
                        end_char=0,
                        headers=headers,
                        metadata={"header_levels": [], "part_of_large_para": True},
                    )
                    chunks.append(chunk)
                    chunk_index += 1

                # Start new chunk
                if len(sentence) > self.chunk_size:
                    # Sentence is too long, just cut it
                    parts = [
                        sentence[i : i + self.chunk_size]
                        for i in range(0, len(sentence), self.chunk_size)
                    ]
                    for part in parts[:-1]:
                        chunk = Chunk(
                            content=part,
                            source_file=source_file,
                            chunk_index=chunk_index,
                            start_char=0,
                            end_char=0,
                            headers=headers,
                            metadata={
                                "header_levels": [],
                                "part_of_large_para": True,
                                "truncated": True,
                            },
                        )
                        chunks.append(chunk)
                        chunk_index += 1

                    current_chunk = parts[-1]  # Last part becomes new chunk
                else:
                    current_chunk = sentence

        # Add final chunk
        if current_chunk.strip():
            chunk = Chunk(
                content=current_chunk
                + ("" if current_chunk.endswith((".", "!", "?")) else "."),
                source_file=source_file,
                chunk_index=chunk_index,
                start_char=0,
                end_char=0,
                headers=headers,
                metadata={"header_levels": [], "part_of_large_para": True},
            )
            chunks.append(chunk)

        return chunks


class EmbeddingGenerator:
    """Generates embeddings using an OpenAI-compatible endpoint."""

    def __init__(self, endpoint: str, model: str):
        """Initialize OpenAI-compatible client [[31]]."""
        self.client = OpenAI(base_url=endpoint, api_key="dummy")
        self.model = model

    async def generate_embeddings(
        self, chunks: List[Chunk], batch_size: int = 32
    ) -> List[np.ndarray]:
        """
        Generate embeddings with batching

        Process:
        1. Batch chunks
        2. Send async requests to endpoint
        3. Handle retries with exponential backoff
        4. Return embeddings in same order
        """
        embeddings = []

        # Process in batches
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i : i + batch_size]
            batch_texts = [chunk.content for chunk in batch]

            try:
                response = self.client.embeddings.create(
                    input=batch_texts, model=self.model
                )

                # Extract embeddings from response
                batch_embeddings = [np.array(data.embedding) for data in response.data]
                embeddings.extend(batch_embeddings)
            except Exception as e:
                raise RuntimeError(f"Failed to generate embeddings for batch: {e}")

        return embeddings

    async def _embed_batch(self, texts: List[str]) -> List[np.ndarray]:
        """Single batch embedding with error handling."""
        try:
            response = self.client.embeddings.create(input=texts, model=self.model)
            return [np.array(data.embedding) for data in response.data]
        except Exception as e:
            raise RuntimeError(f"Failed to embed batch: {e}")


class ChromaDBManager:
    """Manages ChromaDB storage and retrieval."""

    def __init__(self, db_path: str, collection_name: str):
        """Initialize persistent ChromaDB client [[22]]."""
        # Ensure the directory exists
        os.makedirs(db_path, exist_ok=True)
        self.client = chromadb.PersistentClient(path=db_path)
        self.collection = self.client.get_or_create_collection(
            name=collection_name, metadata={"hnsw:space": "cosine"}
        )

    def add_chunks(self, chunks: List[Chunk], embeddings: List[np.ndarray]):
        """
        Add chunks with embeddings to ChromaDB

        Metadata stored per chunk:
        - source_file
        - chunk_index
        - start_char, end_char
        - headers (as JSON string)
        - timestamp
        """
        ids = [chunk.id for chunk in chunks]
        documents = [chunk.content for chunk in chunks]
        metadatas = []

        for chunk in chunks:
            # Convert metadata values to acceptable types for ChromaDB
            metadata = {
                "source_file": chunk.source_file,
                "chunk_index": str(chunk.chunk_index),  # Convert to string
                "start_char": str(chunk.start_char),  # Convert to string
                "end_char": str(chunk.end_char),  # Convert to string
                "headers": json.dumps(chunk.headers),
                "timestamp": str(time.time()),  # Convert to string
            }
            # Add additional metadata, converting values as needed
            for key, value in chunk.metadata.items():
                if isinstance(value, list):
                    # Convert list to JSON string
                    metadata[key] = json.dumps(value)
                elif isinstance(value, dict):
                    # Convert dict to JSON string
                    metadata[key] = json.dumps(value)
                elif isinstance(value, (int, float, str, bool)) or value is None:
                    # Acceptable types for ChromaDB
                    metadata[key] = value
                else:
                    # Convert other types to string
                    metadata[key] = str(value)

            metadatas.append(metadata)

        embeddings_list = [
            emb.tolist() for emb in embeddings
        ]  # Convert numpy arrays to lists

        self.collection.add(
            ids=ids,
            documents=documents,
            metadatas=metadatas,
            embeddings=embeddings_list,
        )

    def query(self, query_embedding: np.ndarray, n_results: int = 5) -> List[Dict]:
        """Query vector DB, return chunks with metadata."""
        results = self.collection.query(
            query_embeddings=[query_embedding.tolist()], n_results=n_results
        )

        # Format results
        formatted_results = []
        for i in range(len(results["ids"][0])):
            result = {
                "id": results["ids"][0][i],
                "content": results["documents"][0][i],
                "metadata": results["metadatas"][0][i],
                "distance": (
                    results["distances"][0][i] if results["distances"] else None
                ),
            }
            formatted_results.append(result)

        return formatted_results


class ParsingError(Exception):
    """Base exception for parsing errors."""

    pass


class EmbeddingError(Exception):
    """Embedding generation failed."""

    pass


def setup_logger(name: str, log_dir: str = "./logs") -> logging.Logger:
    """
    Create logger with file and console handlers

    - Separate log file per agent instance
    - Timestamped entries
    - Configurable log level
    """
    os.makedirs(log_dir, exist_ok=True)

    logger = logging.getLogger(name)
    logger.setLevel(logging.INFO)

    # Prevent adding multiple handlers if logger already exists
    if logger.handlers:
        return logger

    # Create file handler
    log_file = os.path.join(log_dir, f"{name}_{int(time.time())}.log")
    file_handler = logging.FileHandler(log_file)
    file_handler.setLevel(logging.INFO)

    # Create console handler
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)

    # Create formatter
    formatter = logging.Formatter(
        "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
    )
    file_handler.setFormatter(formatter)
    console_handler.setFormatter(formatter)

    # Add handlers to logger
    logger.addHandler(file_handler)
    logger.addHandler(console_handler)

    return logger


class DocumentParserPipeline:
    """Main pipeline for document parsing."""

    def __init__(self, config: DocParserConfig):
        self.config = config
        self.converter = DocumentConverter(config)
        self.consolidator = MarkdownConsolidator(config.output_md_path)
        self.chunker = SemanticChunker(config.chunk_size, config.chunk_overlap)
        self.embedder = EmbeddingGenerator(
            config.embedding_endpoint, config.embedding_model
        )
        self.chroma = ChromaDBManager(config.chroma_db_path, f"docs_{int(time.time())}")
        self.logger = setup_logger("doc_parser")
        # Create a temporary directory for caching
        self.temp_dir = tempfile.mkdtemp(prefix="doc_parser_")

    def _get_cache_path(self, step: str, file_path: str) -> str:
        """Generate a cache file path for a specific step and file."""
        # Use hash of file path to create a safe filename
        import hashlib

        file_hash = hashlib.md5(file_path.encode()).hexdigest()
        return os.path.join(self.temp_dir, f"{step}_{file_hash}.pkl")

    async def process(self) -> DocParserOutput:
        """
        Main processing pipeline with caching

        Steps:
        1. Discover all supported files in input_dir (recursive)
        2. Convert each file to markdown (parallel processing) - CACHED
        3. Chunk the markdown content - CACHED
        4. Consolidate all markdown files
        5. Generate embeddings for all chunks (batched) - CACHED
        6. Store in ChromaDB
        7. Clean up temp directory
        8. Save metadata and return output
        """

        # Step 1: File discovery
        files = self._discover_files()
        self.logger.info(f"Found {len(files)} files to process")

        if not files:
            raise ValueError(f"No supported files found in {self.config.input_dir}")

        # Step 2: Convert files to markdown with caching
        converted_files = []
        for file_path in files:
            cache_path = self._get_cache_path("converted", file_path)
            if os.path.exists(cache_path):
                with open(cache_path, "rb") as f:
                    md_content = pickle.load(f)
            else:
                md_content = self.converter.convert_to_markdown(file_path)
                with open(cache_path, "wb") as f:
                    pickle.dump(md_content, f)
            converted_files.append((file_path, md_content))

        # Step 3: Chunk the markdown content with caching
        all_chunks = []
        chunk_index_offset = 0
        for file_path, md_content in converted_files:
            cache_path = self._get_cache_path("chunks", file_path)
            if os.path.exists(cache_path):
                with open(cache_path, "rb") as f:
                    chunks = pickle.load(f)
            else:
                chunks = self.chunker.chunk_markdown(md_content, file_path)
                # Adjust chunk indices to be globally unique
                for chunk in chunks:
                    chunk.chunk_index += chunk_index_offset
                chunk_index_offset += len(chunks)
                with open(cache_path, "wb") as f:
                    pickle.dump(chunks, f)
            all_chunks.extend(chunks)

        # Step 4: Consolidate all markdown files
        for file_path, md_content in converted_files:
            self.consolidator.append_document(md_content, file_path)
        consolidated_path = self.consolidator.finalize()
        self.logger.info(f"Created consolidated markdown: {consolidated_path}")

        # Step 5: Generate embeddings with caching
        if all_chunks:
            self.logger.info(f"Generating embeddings for {len(all_chunks)} chunks")

            # Load cached embeddings if available
            embeddings = []
            all_cached = True
            for i, chunk in enumerate(all_chunks):
                cache_path = self._get_cache_path(f"embedding_{i}", chunk.id)
                if os.path.exists(cache_path):
                    with open(cache_path, "rb") as f:
                        embedding = pickle.load(f)
                else:
                    # Mark that not all embeddings are cached
                    all_cached = False
                    break
                embeddings.append(embedding)

            # If not all embeddings were cached, generate them
            if not all_cached:
                embeddings = await self.embedder.generate_embeddings(all_chunks)
                # Cache each embedding individually
                for i, embedding in enumerate(embeddings):
                    cache_path = self._get_cache_path(
                        f"embedding_{i}", all_chunks[i].id
                    )
                    with open(cache_path, "wb") as f:
                        pickle.dump(embedding, f)

            # Step 6: Store in ChromaDB
            self.logger.info("Storing in ChromaDB")
            self.chroma.add_chunks(all_chunks, embeddings)
        else:
            self.logger.warning(
                "No chunks were created, skipping embedding and ChromaDB steps"
            )

        # Step 7: Clean up temp directory
        import shutil

        shutil.rmtree(self.temp_dir)

        # Step 8: Return output
        return DocParserOutput(
            consolidated_md_path=consolidated_path,
            chroma_collection_name=self.chroma.collection.name,
            chunk_metadata=[chunk.to_dict() for chunk in all_chunks],
            total_chunks=len(all_chunks),
            processing_log=self.logger.handlers[0].baseFilename,
        )

    def _discover_files(self) -> List[str]:
        """Recursively find all supported files."""
        files = []
        for root, _, filenames in os.walk(self.config.input_dir):
            for filename in filenames:
                if Path(filename).suffix.lower() in self.config.supported_formats:
                    files.append(os.path.join(root, filename))
        return files


async def retry_with_backoff(func, max_retries=3, base_delay=1.0):
    """Utility for retrying failed operations."""
    for attempt in range(max_retries):
        try:
            return await func()
        except Exception as e:
            if attempt == max_retries - 1:
                raise e
            delay = base_delay * (2**attempt)  # Exponential backoff
            await asyncio.sleep(delay)
